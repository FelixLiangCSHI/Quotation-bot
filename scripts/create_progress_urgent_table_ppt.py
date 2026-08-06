from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPT_DIR = ROOT / "PPT"
RULES_DIR = ROOT / "rules"
TEMPLATE_PATH = PPT_DIR / "2026 CSHI Corporate Template - Public Use.pptx"
OUTPUT_PATH = PPT_DIR / "Quotation_Bot_Progress_Urgent_Table_EN.pptx"

FONT = "Noto Sans SC"
TITLE = RGBColor(0x1F, 0x29, 0x37)
BODY = RGBColor(0x4B, 0x55, 0x63)
ORANGE = RGBColor(0xF5, 0x80, 0x23)
DEEP_ORANGE = RGBColor(0xC2, 0x41, 0x0C)
LIGHT_ORANGE = RGBColor(0xFF, 0xF3, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xF8, 0xFA, 0xFC)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
RED = RGBColor(0xDC, 0x26, 0x26)


def main() -> None:
    summary = read_json(RULES_DIR / "rule_review_summary.json")
    merged = read_json(RULES_DIR / "merged_rules.json")

    deck = Presentation(str(TEMPLATE_PATH))
    template_layout = deck.slides[6].slide_layout
    reset_to_layout(deck, template_layout, 5)
    slides = list(deck.slides)

    add_title(slides[0], summary, merged)
    add_visual_summary_table(slides[1], summary, merged)
    add_urgent_table(slides[2], summary)
    add_review_breakdown(slides[3], summary)
    add_next_actions(slides[4], summary)
    add_page_numbers(deck)

    deck.save(str(OUTPUT_PATH))
    print(OUTPUT_PATH)


def read_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def reset_to_layout(deck: Presentation, layout, count: int) -> None:
    slide_id_list = deck.slides._sldIdLst
    for slide_id in list(slide_id_list):
        relationship_id = slide_id.rId
        deck.part.drop_rel(relationship_id)
        slide_id_list.remove(slide_id)
    for _ in range(count):
        slide = deck.slides.add_slide(layout)
        for shape in list(slide.shapes):
            shape._element.getparent().remove(shape._element)


def add_title(slide, summary: dict, merged: dict) -> None:
    add_rect(slide, 0, 0, 0.25, 7.5, ORANGE, ORANGE)
    add_text(slide, "Quotation Bot", 0.85, 0.85, 7.8, 0.55, 34, ORANGE, bold=True)
    add_text(slide, "Progress and Urgent Items", 0.85, 1.55, 8.8, 0.8, 42, TITLE, bold=True)
    add_text(slide, "Based on current JSON rule artifacts and implemented rule engine status.", 0.88, 2.65, 9.8, 0.38, 24, BODY)
    metrics = [
        (str(summary["confirmed_rule_count"]), "confirmed rules"),
        (str(summary["review_rule_count"]), "rules needing review"),
        (str(merged.get("human_approved_rule_count", 0)), "human approved rules"),
        ("18", "unit tests pass"),
    ]
    add_metric_cards(slide, metrics, 0.85, 3.55)
    add_callout(slide, "Core message: rule engine is implemented; urgent blockers are SME review, API exposure, Azure OpenAI access, workflow confirmation, and UAT cases.")


def add_visual_summary_table(slide, summary: dict, merged: dict) -> None:
    add_header(slide, "What Has Been Completed")
    rows = [
        ["Data foundation", "Snapshot data loaded", "380 products / 984 rule signals", "Done"],
        ["Rule engine", "Core validation paths implemented", "Region, system, detector/grid, generator/tube", "Done"],
        ["Confirmed rules", "Executable rule artifact generated", f"{summary['confirmed_rule_count']} confirmed rules", "Done"],
        ["Rule categories", "Structured rule categories counted", "18 region / 590 system / 33 grid / 59 tube", "Done"],
        ["Merged artifact", "Combined rule file exists", f"human approved rules: {merged.get('human_approved_rule_count', 0)}", "Ready"],
        ["Review workflow", "Candidate rules separated", f"{summary['review_rule_count']} need SME review", "Needs action"],
        ["Tests", "Regression coverage exists", "18 unit tests pass", "Done"],
    ]
    add_table(slide, ["Area", "Completed work", "Evidence", "Status"], rows, 0.4, 1.18, 12.55, 5.65, 14)


def add_urgent_table(slide, summary: dict) -> None:
    add_header(slide, "Urgent Items to Resolve")
    review_counts = summary["review_counts_by_type"]
    rows = [
        ["1", "SME rule review", f"{summary['review_rule_count']} rules need confirmation", "PLM / SME / regional specialists"],
        ["2", "Free-text normalization", f"{review_counts['free_text_constraint']} free-text constraints", "SME + rule owner"],
        ["3", "Detector / bucky mapping", f"{review_counts['detector_bucky_match']} detector_bucky_match rules", "Product specialist + PLM"],
        ["4", "Region exclusion decision", f"{review_counts['region_exclusion']} region_exclusion rules", "PLM / regional owner"],
        ["5", "Must-select / choose-one logic", f"{review_counts['must_select']} must_select + {review_counts['any_one_of_n']} any_one_of_n", "PLM + sales support"],
        ["6", "Validation service exposure", "Rule engine not yet exposed as service", "Developer / IT API owner"],
        ["7", "Azure OpenAI access", "Endpoint, deployment, auth, quota, data scope", "IT / AI platform / security"],
        ["8", "Workflow field confirmation", "Required fields for FMT/OTC validation", "PLM + sales support"],
        ["9", "UAT examples", "10-20 real quote cases", "Sales / PLM / SME"],
    ]
    add_table(slide, ["Priority", "Urgent item", "Current gap", "Owner needed"], rows, 0.35, 1.05, 12.65, 5.85, 12)


def add_review_breakdown(slide, summary: dict) -> None:
    add_header(slide, "Rules Needing Review - Breakdown")
    counts = summary["review_counts_by_type"]
    rows = [
        ["free_text_constraint", str(counts["free_text_constraint"]), "Normalize or mark as info-only"],
        ["detector_bucky_match", str(counts["detector_bucky_match"]), "Confirm detector/bucky fields"],
        ["detector_grid_match", str(counts["detector_grid_match"]), "Confirm matrix behavior"],
        ["region_exclusion", str(counts["region_exclusion"]), "Confirm blocked/allowed regions"],
        ["feature_requirement", str(counts["feature_requirement"]), "Convert to when/then requirement"],
        ["must_select", str(counts["must_select"]), "Confirm required scope"],
        ["any_one_of_n", str(counts["any_one_of_n"]), "Confirm min/max selection rule"],
        ["region_only", str(counts["region_only"]), "Confirm region and extra text"],
    ]
    add_table(slide, ["Review-needed type", "Count", "Required action"], rows, 0.65, 1.05, 7.4, 5.8, 14)
    add_rect(slide, 8.35, 1.1, 4.1, 5.75, LIGHT_ORANGE, ORANGE)
    add_text(slide, "Top blocker", 8.7, 1.45, 3.2, 0.3, 26, DEEP_ORANGE, bold=True)
    add_text(slide, "206", 8.7, 2.1, 2.0, 0.6, 44, RED, bold=True)
    add_text(slide, "free-text constraints", 8.7, 2.85, 3.2, 0.35, 22, TITLE, bold=True)
    add_text(slide, "These should not become hard validation rules until SME review is completed.", 8.7, 3.5, 3.25, 1.15, 20, BODY)
    add_text(slide, "Next action", 8.7, 5.05, 2.4, 0.3, 24, DEEP_ORANGE, bold=True)
    add_text(slide, "Assign PLM / SME reviewers by rule type.", 8.7, 5.52, 3.25, 0.7, 20, TITLE, bold=True)


def add_next_actions(slide, summary: dict) -> None:
    add_header(slide, "Next Actions for This Week")
    rows = [
        ["1", "Confirm Azure OpenAI access", "Endpoint, deployment, API version, auth, quota, data scope", "IT / AI platform"],
        ["2", "Expose validation service", "Wrap existing rule engine as callable backend", "Developer / IT API owner"],
        ["3", "Assign SME review owners", f"{summary['review_rule_count']} candidate rules", "PLM / regional specialists"],
        ["4", "Confirm workflow fields", "Required fields for region, system, detector/grid, generator/tube", "PLM / sales support"],
        ["5", "Collect real UAT cases", "10-20 quote questions and expected answers", "Sales / PLM / SME"],
    ]
    add_table(slide, ["Step", "Next action", "Output", "Owner"], rows, 0.45, 1.25, 12.45, 4.25, 14)
    add_callout(slide, "Reporting message: backend progress is real; current urgency is ownership, review, API exposure, and test cases.")


def add_metric_cards(slide, metrics: list[tuple[str, str]], left: float, top: float) -> None:
    for index, (value, label) in enumerate(metrics):
        x = left + index * 3.0
        add_rect(slide, x, top, 2.55, 1.35, LIGHT_ORANGE if index == 1 else WHITE, ORANGE)
        add_text(slide, value, x + 0.22, top + 0.2, 1.4, 0.45, 34, RED if index == 1 else ORANGE, bold=True)
        add_text(slide, label, x + 0.22, top + 0.82, 2.0, 0.25, 18, TITLE, bold=True)


def add_header(slide, title: str) -> None:
    add_text(slide, title, 0.55, 0.32, 12.0, 0.55, 36, TITLE, bold=True)
    add_line(slide, 0.58, 1.02, 1.5, ORANGE)


def add_table(slide, headers: list[str], rows: list[list[str]], left: float, top: float, width: float, height: float, font_size: int) -> None:
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table
    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ORANGE
        format_cell(cell, WHITE, font_size, True)
    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_index % 2 else CARD
            format_cell(cell, BODY, font_size, False)


def add_rect(slide, left: float, top: float, width: float, height: float, fill: RGBColor, line: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)


def add_line(slide, left: float, top: float, length: float, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(length), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def add_callout(slide, text: str) -> None:
    add_rect(slide, 0.55, 6.25, 12.25, 0.55, LIGHT_ORANGE, ORANGE)
    add_text(slide, text, 0.85, 6.42, 11.5, 0.2, 18, TITLE, bold=True)


def add_text(slide, text: str, left: float, top: float, width: float, height: float, size: int, color: RGBColor, bold: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.line_spacing = 1.0


def format_cell(cell, color: RGBColor, size: int, bold: bool) -> None:
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color


def add_page_numbers(deck: Presentation) -> None:
    total = len(deck.slides)
    for index, slide in enumerate(deck.slides, start=1):
        add_text(slide, f"p.{index}/{total}", 11.78, 7.05, 0.9, 0.16, 10, BODY)


if __name__ == "__main__":
    main()