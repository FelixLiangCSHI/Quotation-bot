from __future__ import annotations

import base64
import os
import re
import urllib.request
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPT_DIR = ROOT / "PPT"
MARKDOWN_PATH = ROOT / "docs" / "quotation_bot_mvp_vs_production_architecture.md"
TEMPLATE_PATH = PPT_DIR / "2026 CSHI Corporate Template - Public Use.pptx"
OUTPUT_PATH = PPT_DIR / "Quotation_Bot_Roadmap_From_Markdown_MermaidImages_EN.pptx"
MERMAID_IMAGE_DIR = PPT_DIR / "generated_mermaid"

FONT = "Noto Sans SC"
TITLE = RGBColor(0x1F, 0x29, 0x37)
BODY = RGBColor(0x4B, 0x55, 0x63)
ORANGE = RGBColor(0xF5, 0x80, 0x23)
DEEP_ORANGE = RGBColor(0xC2, 0x41, 0x0C)
LIGHT_ORANGE = RGBColor(0xFF, 0xF3, 0xE8)
BLUE_LIGHT = RGBColor(0xE0, 0xF2, 0xFE)
GREEN_LIGHT = RGBColor(0xDC, 0xFC, 0xE7)
GRAY_LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)


PHASES = [
    {
        "title": "Phase 0 - Align Scope and Access",
        "conclusion": "Confirm the pilot scenario, Azure OpenAI access, and data usage boundary before building more UI.",
        "what": [
            "Define the first pilot scenario.",
            "Confirm Azure OpenAI endpoint and model access.",
            "Confirm whether quotation/product/rule data can be sent to Azure OpenAI.",
            "Confirm whether the first demo runs locally, internally, or behind an API gateway.",
        ],
        "support": [
            "IT / AI platform team: Azure OpenAI endpoint, deployment name, API version.",
            "IT / security / data owner: data usage approval.",
            "Mentor, PLM, Sales / BDM: pilot business scenario.",
            "Sales / PLM: first 10-20 real cases.",
        ],
        "deliverable": "One-page pilot scope, approved model/API access path, and initial test question set.",
    },
    {
        "title": "Phase 1 - Build Beta Version Chat Interface",
        "conclusion": "Use Streamlit as frontend first.",
        "what": [
            "Build a local chat UI.",
            "Allow users to input quote/configuration questions.",
            "Display extracted fields, validation result, and explanation.",
            "Use session state only for the current conversation.",
        ],
        "support": [
            "Mentor / sponsor: approval to run internal demo.",
            "IT / AI platform: Azure OpenAI credentials or approved access method.",
            "Sales / PLM: demo questions.",
        ],
        "deliverable": "Working Beta version UI and demo flow: user question -> extracted fields -> validation -> explanation.",
    },
    {
        "title": "Phase 2 - Connect Reasoning Layer",
        "conclusion": "Use Azure OpenAI as the preferred reasoning layer. Use open-source LLM only if Azure OpenAI is unavailable or policy requires self-hosting.",
        "what": [
            "Use Azure OpenAI for intent extraction.",
            "Use Azure OpenAI for field extraction.",
            "Use Azure OpenAI for explanation wording.",
            "Do not use Azure OpenAI as the final validation authority.",
        ],
        "support": [
            "IT / AI platform: AZURE_OPENAI_ENDPOINT, deployment name, API version.",
            "IT / security: authentication method.",
            "Security / data owner: allowed data scope.",
            "IT / AI platform: quota / rate limit.",
        ],
        "deliverable": "parse_intent_and_fields function or service, structured extraction prompt, and explanation prompt using rule-engine output only.",
    },
    {
        "title": "Phase 3 - Use Existing Rule Engine as Validation Authority",
        "conclusion": "The existing rule engine should remain the validation authority. The chatbot should call it, not recreate validation inside the LLM.",
        "what": [
            "Keep using QuotationRuleEngine for deterministic validation.",
            "Load quotation_snapshot.json and rules/merged_rules.json.",
            "Convert parsed user input into structured validation input.",
            "Return valid, invalid, incomplete, warning, and info results.",
        ],
        "support": [
            "Developer / intern: rule engine API wrapper.",
            "PLM / SME: business validation of messages.",
            "QA / Sales / PLM: regression test cases.",
        ],
        "deliverable": "Validation service around the existing rule engine. First endpoint: POST /validation/check.",
    },
    {
        "title": "Phase 4 - Keep Beta Version Data Simple",
        "conclusion": "For Beta version, read data directly from JSON and Markdown.",
        "what": [
            "Continue using quotation_snapshot.json as product/data source.",
            "Continue using rules/merged_rules.json as rule artifact.",
            "Use Markdown docs for implementation notes and workflow explanation.",
            "Add database/search only when pilot needs multi-user access, refresh control, audit, or performance.",
        ],
        "support": [
            "Data owner: source file ownership and refresh decision.",
            "Developer: file-based data loader.",
            "IT / security: production data storage decision later.",
        ],
        "deliverable": "File-based Beta version data loader and decision that database/search/embedding are later-stage options.",
    },
    {
        "title": "Phase 5 - Add Minimal Memory",
        "conclusion": "The bot needs short-term session memory for multi-turn conversations. Long-term memory is not required for Beta version.",
        "what": [
            "Store current quote fields during a session.",
            "Remember missing fields already asked.",
            "Reset or export session after the demo.",
            "Avoid storing long-term chat history until IT/security approves retention policy.",
        ],
        "support": [
            "IT / security: chat history retention decision.",
            "Developer: session schema.",
            "IT / security: user identity model for production conversation ownership.",
        ],
        "deliverable": "Beta version session state in Streamlit/Gradio or local Python object; production memory recommendation later.",
    },
    {
        "title": "Phase 6 - Expose Company APIs",
        "conclusion": "After Beta version works locally, convert useful functions into company-owned APIs for future frontends.",
        "what": [
            "Build validation service first.",
            "Build product search API.",
            "Build quote context API.",
            "Build workflow required-fields API.",
            "Add review case API after SME workflow is agreed.",
        ],
        "support": [
            "Rule engine developer: validation service.",
            "Data/API developer: product APIs.",
            "Backend developer: quote context API.",
            "PLM + developer: workflow required-fields API.",
            "SME + developer: review case API.",
        ],
        "deliverable": "Internal API contract, endpoint owner list, and first callable validation endpoint.",
    },
    {
        "title": "Phase 7 - Decide Production Frontend",
        "conclusion": "Frontend should be replaceable.",
        "what": [
            "Use Streamlit for Beta version.",
            "After APIs are ready, evaluate Teams, internal Web App, Dify, Coze, or Copilot Studio.",
            "Choose production frontend based on company security, licensing, deployment, and user adoption.",
        ],
        "support": [
            "Sponsor: production direction.",
            "IT: approved platform and hosting model.",
            "Users: adoption feedback.",
            "Security: third-party platform approval if Coze/Dify is considered.",
        ],
        "deliverable": "Production frontend decision after API contract is reviewed.",
    },
    {
        "title": "Phase 8 - SME Rule Review and Expansion",
        "conclusion": "The rule engine is implemented, but not every extracted rule should be automated immediately. The 387 review-needed rules require SME ownership.",
        "what": [
            "Assign SME owners for review-needed rule categories.",
            "Review free-text constraints and region exclusions first.",
            "Normalize approved rules into structured payloads.",
            "Add rule engine handlers for approved types.",
        ],
        "support": [
            "PLM / SME: rule review owner.",
            "Regional product specialists: regional validation.",
            "Developer: handler implementation.",
            "QA / developer: regression tests.",
        ],
        "deliverable": "Reviewed rules, expanded rule handlers, and updated rules/merged_rules.json.",
    },
    {
        "title": "Phase 9 - Pilot and UAT",
        "conclusion": "Pilot should validate a few high-value quote scenarios, not the entire product universe.",
        "what": [
            "Select 10-20 real quote questions.",
            "Run them through the Beta version.",
            "Compare bot output with SME expectation.",
            "Record false positives, false negatives, missing fields, and unclear explanations.",
            "Improve prompts, workflow fields, and rule handlers.",
        ],
        "support": [
            "Sales / PLM: real quote cases.",
            "SME / PLM: expected answers.",
            "QA / project owner: UAT tracker.",
            "Developer / AI engineer: improvement owner.",
        ],
        "deliverable": "UAT result table and go/no-go recommendation for pilot expansion.",
    },
]


def main() -> None:
    mermaid_images = render_mermaid_images()
    deck = Presentation(str(TEMPLATE_PATH))
    template_layout = deck.slides[6].slide_layout
    reset_to_layout(deck, template_layout, 16)

    slides = list(deck.slides)
    add_title(slides[0])
    add_executive(slides[1])
    add_beta_mermaid(slides[2], mermaid_images[0])
    add_production_mermaid(slides[3], mermaid_images[1])
    for index, phase in enumerate(PHASES, start=4):
        add_phase_slide(slides[index], phase)
    add_roadmap(slides[14])
    add_final(slides[15])
    add_page_numbers(deck)

    deck.save(str(OUTPUT_PATH))
    print(OUTPUT_PATH)


def render_mermaid_images() -> list[Path]:
    if os.getenv("ALLOW_EXTERNAL_MERMAID_RENDER", "").strip().casefold() not in ("1", "true", "yes"):
        raise RuntimeError(
            "Rendering Mermaid diagrams sends internal architecture content to "
            "the external mermaid.ink service. Set ALLOW_EXTERNAL_MERMAID_RENDER=1 "
            "to explicitly opt in, or render the diagrams locally."
        )
    markdown = MARKDOWN_PATH.read_text(encoding="utf-8")
    blocks = re.findall(r"```mermaid\s*(.*?)```", markdown, flags=re.S)
    if len(blocks) < 2:
        raise ValueError("Expected at least two Mermaid blocks in the roadmap markdown.")

    MERMAID_IMAGE_DIR.mkdir(exist_ok=True)
    image_paths: list[Path] = []
    for index, block in enumerate(blocks[:2], start=1):
        encoded = base64.urlsafe_b64encode(block.strip().encode("utf-8")).decode("ascii").rstrip("=")
        url = f"https://mermaid.ink/img/{encoded}?type=png&bgColor=white"
        image_path = MERMAID_IMAGE_DIR / f"roadmap_mermaid_{index}.png"
        request = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(request, timeout=30) as response:
            image_path.write_bytes(response.read())
        image_paths.append(image_path)
    return image_paths


def reset_to_layout(deck: Presentation, layout, count: int) -> None:
    slide_id_list = deck.slides._sldIdLst
    for slide_id in list(slide_id_list):
        relationship_id = slide_id.rId
        deck.part.drop_rel(relationship_id)
        slide_id_list.remove(slide_id)
    for _ in range(count):
        slide = deck.slides.add_slide(layout)
        for shape in list(slide.shapes):
            shape._element.getparent().remove(shape._element)


def add_title(slide) -> None:
    add_rect(slide, 0.0, 0.0, 0.25, 7.5, ORANGE, ORANGE)
    add_text(slide, "Quotation Bot", 0.85, 0.82, 7.8, 0.55, 36, ORANGE, bold=True)
    add_text(slide, "Implementation Roadmap\nFrom Markdown", 0.85, 1.55, 9.0, 1.55, 44, TITLE, bold=True)
    add_text(slide, "Beta version and Production roadmap with Mermaid architecture pages and phase-by-phase execution pages.", 0.88, 3.55, 9.8, 0.7, 24, BODY)
    add_badge_row(slide, 0.88, 5.05, ["Mermaid diagrams", "Conclusion first", "What to do", "Required support", "Deliverable"])
    add_text(slide, "Meeting draft | 2026-07-07", 0.88, 6.25, 5.2, 0.3, 20, BODY)


def add_executive(slide) -> None:
    add_header(slide, "Executive Conclusion")
    add_rect(slide, 0.55, 1.2, 12.25, 1.55, LIGHT_ORANGE, ORANGE)
    add_text(slide, "Recommended path: Beta version first, Production later.", 0.9, 1.62, 11.0, 0.35, 28, TITLE, bold=True)
    add_two_columns(
        slide,
        "Beta version",
        [
            "Lightweight frontend.",
            "Azure OpenAI for extraction and explanation.",
            "Existing JSON / Markdown files.",
            "Existing rule engine as validation authority.",
        ],
        "Production",
        [
            "Enterprise frontend and authentication.",
            "Internal APIs and persistent memory.",
            "Monitoring and SME review workflow.",
            "Optional search, embedding, and database.",
        ],
        0.55,
        3.0,
    )
    add_deliverable(slide, "Current asset: Existing QuotationRuleEngine. Next action: expose it as a callable validation service.")


def add_beta_mermaid(slide, image_path: Path) -> None:
    add_header(slide, "Beta Version Architecture - Mermaid Flow")
    add_text(slide, "Generated directly from the first Mermaid block in the roadmap markdown.", 0.65, 1.05, 10.5, 0.28, 20, BODY)
    add_picture_fit(slide, image_path, 0.35, 1.38, 12.65, 5.45)


def add_production_mermaid(slide, image_path: Path) -> None:
    add_header(slide, "Production Architecture - Mermaid Flow")
    add_text(slide, "Generated directly from the second Mermaid block in the roadmap markdown.", 0.65, 1.05, 10.5, 0.28, 20, BODY)
    add_picture_fit(slide, image_path, 0.35, 1.38, 12.65, 5.45)


def add_phase_slide(slide, phase: dict[str, object]) -> None:
    add_header(slide, phase["title"])
    add_rect(slide, 0.55, 1.18, 12.25, 1.05, LIGHT_ORANGE, ORANGE)
    add_text(slide, "Conclusion", 0.85, 1.36, 1.8, 0.25, 22, DEEP_ORANGE, bold=True)
    add_text(slide, str(phase["conclusion"]), 2.55, 1.32, 9.7, 0.38, 20, TITLE, bold=True)
    add_two_columns(slide, "What to do", phase["what"], "Required support", phase["support"], 0.55, 2.55)
    add_deliverable(slide, str(phase["deliverable"]))


def add_roadmap(slide) -> None:
    add_header(slide, "Final Recommended Roadmap")
    rows = [[str(i), p["title"].replace(f"Phase {i} - ", ""), "; ".join(p["support"][:2]), str(p["deliverable"])] for i, p in enumerate(PHASES)]
    add_table(slide, ["Phase", "What to do", "Required support", "Deliverable"], rows, 0.35, 1.25, 12.65, 5.7, 11)


def add_final(slide) -> None:
    add_header(slide, "Final Recommendation")
    add_numbered_list(
        slide,
        [
            "Use Azure OpenAI for intent extraction and explanation.",
            "Use Streamlit or Gradio first for the Beta version frontend.",
            "Read JSON and Markdown directly for Beta version.",
            "Keep the existing rule engine as the validation authority.",
            "Expose POST /validation/check as the first internal API.",
            "Add database, embeddings, search index, and formal frontend after Beta version proves value.",
        ],
        0.7,
        1.35,
        11.6,
        4.6,
        23,
    )
    add_deliverable(slide, "Final principle: Frontend can change. Reasoning provider can change. Data store can evolve. The existing rule engine remains the deterministic validation authority.")


def add_header(slide, title: str) -> None:
    add_text(slide, title, 0.55, 0.32, 12.0, 0.55, 36, TITLE, bold=True)
    add_line(slide, 0.58, 1.02, 1.5, ORANGE)


def add_two_columns(slide, left_title: str, left_items: list[str], right_title: str, right_items: list[str], left: float, top: float) -> None:
    add_rect(slide, left, top, 5.95, 3.65, WHITE, BORDER)
    add_text(slide, left_title, left + 0.28, top + 0.25, 4.7, 0.32, 24, ORANGE, bold=True)
    add_bullets(slide, left_items, left + 0.28, top + 0.78, 5.35, 2.45, 16)
    add_rect(slide, left + 6.25, top, 5.95, 3.65, LIGHT_ORANGE, ORANGE)
    add_text(slide, right_title, left + 6.53, top + 0.25, 4.8, 0.32, 24, DEEP_ORANGE, bold=True)
    add_bullets(slide, right_items, left + 6.53, top + 0.78, 5.35, 2.45, 15)


def add_deliverable(slide, text: str) -> None:
    add_rect(slide, 0.55, 6.3, 12.25, 0.58, LIGHT_ORANGE, ORANGE)
    add_text(slide, "Deliverable", 0.82, 6.48, 1.65, 0.2, 18, DEEP_ORANGE, bold=True)
    add_text(slide, text, 2.45, 6.45, 10.0, 0.24, 16, TITLE, bold=True)


def flow_box(slide, text: str, left: float, top: float, width: float, height: float, fill: RGBColor = WHITE) -> None:
    add_rect(slide, left, top, width, height, fill, ORANGE)
    add_center_text(slide, text, left + 0.08, top + height / 2 - 0.18, width - 0.16, 0.3, 14, TITLE, bold=True)


def add_picture_fit(slide, image_path: Path, left: float, top: float, width: float, height: float) -> None:
    with Image.open(image_path) as image:
        image_width, image_height = image.size
    image_ratio = image_width / image_height
    box_ratio = width / height
    if image_ratio >= box_ratio:
        final_width = width
        final_height = width / image_ratio
    else:
        final_height = height
        final_width = height * image_ratio
    final_left = left + (width - final_width) / 2
    final_top = top + (height - final_height) / 2
    slide.shapes.add_picture(str(image_path), Inches(final_left), Inches(final_top), width=Inches(final_width), height=Inches(final_height))


def draw_h_arrow(slide, x1: float, y: float, x2: float) -> None:
    add_line(slide, x1, y, max(0.05, x2 - x1), ORANGE)
    add_text(slide, ">", x2 - 0.05, y - 0.11, 0.16, 0.16, 16, ORANGE, bold=True)


def draw_v_arrow(slide, x: float, y1: float, y2: float) -> None:
    add_line(slide, x, y1, max(0.05, y2 - y1), ORANGE, vertical=True)
    add_text(slide, "v", x - 0.07, y2 - 0.08, 0.16, 0.16, 16, ORANGE, bold=True)


def add_table(slide, headers: list[str], rows: list[list[str]], left: float, top: float, width: float, height: float, font_size: int) -> None:
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table
    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ORANGE
        format_cell(cell, WHITE, font_size, True)
    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_index % 2 else GRAY_LIGHT
            format_cell(cell, BODY, font_size, False)


def add_numbered_list(slide, items: list[str], left: float, top: float, width: float, height: float, size: int) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    for index, item in enumerate(items, start=1):
        paragraph = frame.paragraphs[0] if index == 1 else frame.add_paragraph()
        paragraph.text = f"{index}. {item}"
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = BODY
        paragraph.space_after = Pt(6)


def add_bullets(slide, bullets: list[str], left: float, top: float, width: float, height: float, size: int) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = BODY
        paragraph.space_after = Pt(4)


def add_badge_row(slide, left: float, top: float, labels: list[str]) -> None:
    x = left
    for label in labels:
        width = max(1.6, len(label) * 0.1)
        add_rect(slide, x, top, width, 0.38, LIGHT_ORANGE, ORANGE)
        add_center_text(slide, label, x + 0.05, top + 0.08, width - 0.1, 0.16, 14, TITLE, bold=True)
        x += width + 0.22


def add_text(slide, text: str, left: float, top: float, width: float, height: float, size: int, color: RGBColor, bold: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.line_spacing = 1.0


def add_center_text(slide, text: str, left: float, top: float, width: float, height: float, size: int, color: RGBColor, bold: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def add_rect(slide, left: float, top: float, width: float, height: float, fill: RGBColor, line: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)


def add_line(slide, left: float, top: float, length: float, color: RGBColor, vertical: bool = False) -> None:
    if vertical:
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(0.03), Inches(length))
    else:
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(length), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def format_cell(cell, color: RGBColor, size: int, bold: bool) -> None:
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color


def add_page_numbers(deck: Presentation) -> None:
    total = len(deck.slides)
    for index, slide in enumerate(deck.slides, start=1):
        add_text(slide, f"p.{index}/{total}", 11.78, 7.05, 0.9, 0.16, 10, BODY)


if __name__ == "__main__":
    main()