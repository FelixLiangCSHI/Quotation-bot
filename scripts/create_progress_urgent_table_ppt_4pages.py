from __future__ import annotations

import csv
import json
from collections import OrderedDict
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPT_DIR = ROOT / "PPT"
RULES_DIR = ROOT / "rules"
TEMPLATE_PATH = PPT_DIR / "2026 CSHI Corporate Template - Public Use.pptx"
OUTPUT_PATH = PPT_DIR / "Quotation_Bot_Progress_Urgent_Table_4Pages_With_Examples_EN.pptx"

FONT = "Noto Sans SC"
TITLE = RGBColor(0x1F, 0x29, 0x37)
BODY = RGBColor(0x4B, 0x55, 0x63)
ORANGE = RGBColor(0xF5, 0x80, 0x23)
DEEP_ORANGE = RGBColor(0xC2, 0x41, 0x0C)
LIGHT_ORANGE = RGBColor(0xFF, 0xF3, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xF8, 0xFA, 0xFC)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
RED = RGBColor(0xDC, 0x26, 0x26)


def main() -> None:
    summary = read_json(RULES_DIR / "rule_review_summary.json")
    merged = read_json(RULES_DIR / "merged_rules.json")
    examples = read_examples()

    deck = Presentation(str(TEMPLATE_PATH))
    layout = deck.slides[6].slide_layout
    reset_to_layout(deck, layout, 4)
    slides = list(deck.slides)

    add_progress_overview(slides[0], summary, merged)
    add_urgent_items(slides[1], summary)
    add_review_examples(slides[2], summary, examples)
    add_category_examples(slides[3], summary, examples)
    add_page_numbers(deck)

    deck.save(str(OUTPUT_PATH))
    print(OUTPUT_PATH)


def read_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def read_examples() -> OrderedDict[str, dict[str, str]]:
    wanted = [
        "free_text_constraint",
        "region_only",
        "any_one_of_n",
        "must_select",
        "detector_bucky_match",
        "detector_grid_match",
        "feature_requirement",
        "region_exclusion",
    ]
    examples: OrderedDict[str, dict[str, str]] = OrderedDict()
    with (RULES_DIR / "rules_needing_review.csv").open("r", encoding="utf-8-sig", newline="") as file:
        for row in csv.DictReader(file):
            rule_type = row["rule_type"]
            if rule_type not in wanted or rule_type in examples:
                continue
            examples[rule_type] = {
                "product_id": row["product_id"] or "-",
                "step_id": row["step_id"] or "-",
                "message": row["message"],
                "source": f"{row['source_sheet']}!{row['source_cell']}",
                "reason": row["review_reason"],
            }
            if len(examples) == len(wanted):
                break
    return examples


def reset_to_layout(deck: Presentation, layout, count: int) -> None:
    slide_id_list = deck.slides._sldIdLst
    for slide_id in list(slide_id_list):
        relationship_id = slide_id.rId
        deck.part.drop_rel(relationship_id)
        slide_id_list.remove(slide_id)
    for _ in range(count):
        slide = deck.slides.add_slide(layout)
        for shape in list(slide.shapes):
            shape._element.getparent().remove(shape._element)


def add_progress_overview(slide, summary: dict, merged: dict) -> None:
    add_header(slide, "Progress Overview - What Has Been Completed")
    metrics = [
        (str(summary["confirmed_rule_count"]), "confirmed rules"),
        (str(summary["review_rule_count"]), "rules needing review"),
        (str(merged.get("human_approved_rule_count", 0)), "human approved rules"),
        ("18", "unit tests pass"),
    ]
    add_metric_cards(slide, metrics, 0.55, 1.15)
    rows = [
        ["Data foundation", "Snapshot and rule data loaded", "380 products / 984 rule signals", "Done"],
        ["Rule engine", "Core validation paths implemented", "Region, system, detector/grid, generator/tube", "Done"],
        ["Confirmed rules", "Executable rule artifact generated", "700 confirmed rules", "Done"],
        ["Rule categories", "Implemented rule categories counted", "18 region / 590 system / 33 grid / 59 tube", "Done"],
        ["Merged artifact", "Combined rule file exists", "0 human approved rules merged", "Ready"],
        ["Review workflow", "Candidate rules separated", "387 need SME review", "Needs action"],
    ]
    add_table(slide, ["Area", "Completed work", "Evidence", "Status"], rows, 0.35, 2.9, 12.65, 3.65, 11)


def add_urgent_items(slide, summary: dict) -> None:
    add_header(slide, "Urgent Items to Resolve")
    counts = summary["review_counts_by_type"]
    rows = [
        ["1", "SME rule review", "387 candidate rules", "PLM / SME / regional specialists"],
        ["2", "Free-text normalization", f"{counts['free_text_constraint']} free-text constraints", "SME + rule owner"],
        ["3", "Detector / bucky mapping", f"{counts['detector_bucky_match']} detector_bucky_match", "Product specialist + PLM"],
        ["4", "Region exclusion decision", f"{counts['region_exclusion']} region_exclusion", "PLM / regional owner"],
        ["5", "Must-select / choose-one logic", f"{counts['must_select']} must_select + {counts['any_one_of_n']} choose-one", "PLM + sales support"],
        ["6", "Validation service exposure", "Rule engine not callable yet", "Developer / IT API owner"],
        ["7", "Azure OpenAI access", "Endpoint, auth, quota, data scope", "IT / AI platform / security"],
        ["8", "Workflow field confirmation", "Required fields for FMT/OTC", "PLM + sales support"],
        ["9", "UAT examples", "10-20 real quote cases", "Sales / PLM / SME"],
    ]
    add_table(slide, ["Priority", "Urgent item", "Current gap", "Owner needed"], rows, 0.3, 1.12, 12.8, 5.75, 10)


def add_review_examples(slide, summary: dict, examples: OrderedDict[str, dict[str, str]]) -> None:
    add_header(slide, "Human Review Cases - Real Examples")
    selected = [
        ("free_text_constraint", examples["free_text_constraint"]),
        ("detector_bucky_match", examples["detector_bucky_match"]),
        ("detector_grid_match", examples["detector_grid_match"]),
        ("region_exclusion", examples["region_exclusion"]),
    ]
    rows = []
    for rule_type, example in selected:
        rows.append([
            rule_type,
            example["product_id"],
            example["source"],
            shorten(example["message"], 95),
            human_action(rule_type),
        ])
    add_table(slide, ["Rule type", "Product", "Source", "Extracted case", "Manual review needed"], rows, 0.28, 1.15, 12.85, 4.0, 9)
    add_rect(slide, 0.55, 5.55, 12.25, 0.85, LIGHT_ORANGE, ORANGE)
    add_text(slide, "Key point", 0.85, 5.78, 1.5, 0.2, 18, DEEP_ORANGE, bold=True)
    add_text(slide, "These cases are already extracted from JSON/CSV, but business meaning still needs SME confirmation before hard-block automation.", 2.15, 5.75, 10.0, 0.24, 16, BODY, bold=True)


def add_category_examples(slide, summary: dict, examples: OrderedDict[str, dict[str, str]]) -> None:
    add_header(slide, "Extracted Categories and Example Cases")
    counts = summary["review_counts_by_type"]
    rows = []
    for rule_type in [
        "free_text_constraint",
        "region_only",
        "any_one_of_n",
        "must_select",
        "detector_bucky_match",
        "detector_grid_match",
        "feature_requirement",
        "region_exclusion",
    ]:
        example = examples[rule_type]
        rows.append([rule_type, str(counts[rule_type]), shorten(example["message"], 80), human_action(rule_type)])
    add_table(slide, ["Extracted category", "Count", "Example", "Next action"], rows, 0.28, 1.05, 12.85, 5.55, 9)


def human_action(rule_type: str) -> str:
    return {
        "free_text_constraint": "Normalize or info-only",
        "region_only": "Confirm region + extra text",
        "any_one_of_n": "Confirm min/max selection",
        "must_select": "Confirm required scope",
        "detector_bucky_match": "Map detector/bucky fields",
        "detector_grid_match": "Confirm grid logic",
        "feature_requirement": "Define when/then rule",
        "region_exclusion": "Confirm blocked regions",
    }[rule_type]


def shorten(text: str, limit: int) -> str:
    text = " ".join(text.split())
    return text if len(text) <= limit else text[: limit - 3] + "..."


def add_metric_cards(slide, metrics: list[tuple[str, str]], left: float, top: float) -> None:
    for index, (value, label) in enumerate(metrics):
        x = left + index * 3.13
        add_rect(slide, x, top, 2.65, 1.2, LIGHT_ORANGE if index == 1 else WHITE, ORANGE)
        add_text(slide, value, x + 0.2, top + 0.18, 1.35, 0.42, 30, RED if index == 1 else ORANGE, bold=True)
        add_text(slide, label, x + 0.2, top + 0.75, 2.1, 0.24, 15, TITLE, bold=True)


def add_header(slide, title: str) -> None:
    add_text(slide, title, 0.5, 0.32, 12.1, 0.5, 34, TITLE, bold=True)
    add_line(slide, 0.54, 0.98, 1.45, ORANGE)


def add_table(slide, headers: list[str], rows: list[list[str]], left: float, top: float, width: float, height: float, font_size: int) -> None:
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table
    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ORANGE
        format_cell(cell, WHITE, font_size, True)
    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_index % 2 else CARD
            format_cell(cell, BODY, font_size, False)


def add_rect(slide, left: float, top: float, width: float, height: float, fill: RGBColor, line: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)


def add_line(slide, left: float, top: float, length: float, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(length), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def add_text(slide, text: str, left: float, top: float, width: float, height: float, size: int, color: RGBColor, bold: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.line_spacing = 1.0


def format_cell(cell, color: RGBColor, size: int, bold: bool) -> None:
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color


def add_page_numbers(deck: Presentation) -> None:
    total = len(deck.slides)
    for index, slide in enumerate(deck.slides, start=1):
        add_text(slide, f"p.{index}/{total}", 11.82, 7.05, 0.9, 0.16, 10, BODY)


if __name__ == "__main__":
    main()