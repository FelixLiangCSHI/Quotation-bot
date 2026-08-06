from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPT_DIR = ROOT / "PPT"
TEMPLATE_PATH = PPT_DIR / "2026 CSHI Corporate Template - Public Use.pptx"
OUTPUT_PATH = PPT_DIR / "Quotation_Bot_API_Flow_Orange_Template7_LargeFont_EN.pptx"

FONT = "Noto Sans SC"
TITLE = RGBColor(0x1F, 0x29, 0x37)
BODY = RGBColor(0x4B, 0x55, 0x63)
ORANGE = RGBColor(0xF5, 0x80, 0x23)
DEEP_ORANGE = RGBColor(0xC2, 0x41, 0x0C)
LIGHT_ORANGE = RGBColor(0xFF, 0xF3, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xF8, 0xFA, 0xFC)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)


def main() -> None:
    deck = Presentation(str(TEMPLATE_PATH))
    template_layout = deck.slides[6].slide_layout
    reset_to_layout(deck, template_layout, 15)

    slides = list(deck.slides)
    add_title(slides[0])
    add_overview(slides[1])
    add_mermaid_flow(slides[2])
    add_reference_architecture(slides[3])
    add_input_session(slides[4])
    add_intent_extraction(slides[5])
    add_workflow_reasoning(slides[6])
    add_grounding(slides[7])
    add_quote_context(slides[8])
    add_validation(slides[9])
    add_explanation(slides[10])
    add_sme_review(slides[11])
    add_endpoint_matrix(slides[12])
    add_collaboration(slides[13])
    add_decisions(slides[14])
    add_page_numbers(deck)

    deck.save(str(OUTPUT_PATH))
    print(OUTPUT_PATH)


def reset_to_layout(deck: Presentation, layout, count: int) -> None:
    slide_id_list = deck.slides._sldIdLst
    for slide_id in list(slide_id_list):
        relationship_id = slide_id.rId
        deck.part.drop_rel(relationship_id)
        slide_id_list.remove(slide_id)
    for _ in range(count):
        slide = deck.slides.add_slide(layout)
        for shape in list(slide.shapes):
            shape._element.getparent().remove(shape._element)


def add_title(slide) -> None:
    add_rect(slide, 0.0, 0.0, 0.25, 7.5, ORANGE, ORANGE)
    add_text(slide, "Quotation Bot", 0.85, 0.82, 7.8, 0.55, 36, ORANGE, bold=True)
    add_text(slide, "API-Driven Chatbot\nImplementation Flow", 0.85, 1.55, 8.4, 1.55, 44, TITLE, bold=True)
    add_text(
        slide,
        "A larger-font, flowchart-first deck converted from the implementation markdown. Focus: urgent endpoints, resource owners, and the practical build path.",
        0.88,
        3.45,
        9.2,
        0.9,
        24,
        BODY,
    )
    add_badge_row(slide, 0.88, 5.05, ["Template 7 Background", "Orange Theme", "Flowcharts + Tables", "Endpoint Focus"])
    add_text(slide, "Meeting draft | 2026-07-07", 0.88, 6.25, 5.2, 0.3, 20, BODY)


def add_overview(slide) -> None:
    add_header(slide, "Overview / Agenda")
    add_rect(slide, 0.62, 1.25, 5.95, 5.55, WHITE, BORDER)
    add_text(slide, "What this deck answers", 0.95, 1.55, 4.9, 0.34, 26, ORANGE, bold=True)
    add_numbered_list(
        slide,
        [
            "What is the end-to-end chatbot flow?",
            "Which APIs and endpoints are urgently needed?",
            "What does each step require from IT, PLM, SME, and data owners?",
            "How does the current rule engine fit into the full solution?",
            "What decisions are needed to start the pilot?",
        ],
        0.95,
        2.1,
        5.0,
        3.7,
        21,
    )

    add_rect(slide, 6.85, 1.25, 5.85, 5.55, LIGHT_ORANGE, ORANGE)
    add_text(slide, "One-line implementation path", 7.18, 1.55, 4.8, 0.34, 26, DEEP_ORANGE, bold=True)
    add_vertical_flow(
        slide,
        7.2,
        2.1,
        [
            "User input",
            "Intent + field extraction",
            "Missing-field reasoning",
            "Product / rule grounding APIs",
            "Quote context",
            "Rule engine validation",
            "Explanation + SME review",
        ],
        box_w=4.8,
        box_h=0.43,
        font_size=20,
    )


def add_mermaid_flow(slide) -> None:
    add_header(slide, "Markdown Flowchart Rendered as PPT Diagram")
    add_text(slide, "Equivalent to the Mermaid flow in the markdown plan", 0.75, 1.08, 8.5, 0.3, 22, BODY)
    nodes = [
        ("Sales User", 0.45, 1.55),
        ("Chat UI", 2.25, 1.55),
        ("Auth + Session", 4.05, 1.55),
        ("Bot Orchestrator", 6.05, 1.55),
        ("Intent Extraction", 8.25, 1.55),
        ("Workflow Resolver", 10.35, 1.55),
        ("Ask Missing Info", 4.25, 3.4),
        ("Product + Rule APIs", 6.45, 3.4),
        ("Quote Context", 8.65, 3.4),
        ("Validation API", 10.65, 3.4),
        ("Explanation", 6.45, 5.25),
        ("SME Review", 8.65, 5.25),
        ("Merged Rules", 10.65, 5.25),
    ]
    for label, x, y in nodes:
        flow_box(slide, label, x, y, 1.65, 0.72, 20)
    arrows = [
        (1.98, 1.9, 2.22, 1.9),
        (3.78, 1.9, 4.02, 1.9),
        (5.78, 1.9, 6.02, 1.9),
        (7.98, 1.9, 8.22, 1.9),
        (9.98, 1.9, 10.32, 1.9),
        (11.2, 2.27, 11.2, 3.36),
        (5.9, 3.75, 6.42, 3.75),
        (8.1, 3.75, 8.62, 3.75),
        (10.3, 3.75, 10.62, 3.75),
        (11.35, 4.13, 11.35, 5.21),
        (10.3, 5.6, 10.62, 5.6),
        (8.1, 5.6, 8.62, 5.6),
    ]
    for x1, y1, x2, y2 in arrows:
        add_arrow(slide, x1, y1, x2, y2)
    add_callout(slide, "Key point: the bot only explains; deterministic endpoints provide facts, state, and validation results.")


def add_reference_architecture(slide) -> None:
    add_header(slide, "Reference Architecture")
    cols = [
        ("Channel", ["Teams / Web", "User auth", "Conversation UI"]),
        ("Orchestrator", ["Session state", "Tool calling", "Routing logic"]),
        ("AI Layer", ["LLM parse", "Prompt templates", "Explanation"]),
        ("Business APIs", ["Product search", "Workflow", "Quote context"]),
        ("Rule Layer", ["Validation", "Merged rules", "SME review"]),
    ]
    x = 0.45
    for title, bullets in cols:
        add_rect(slide, x, 1.45, 2.35, 4.8, WHITE, BORDER)
        add_text(slide, title, x + 0.18, 1.78, 2.0, 0.3, 22, ORANGE, bold=True)
        add_bullets(slide, bullets, x + 0.18, 2.35, 1.95, 2.7, 20)
        x += 2.5
    add_callout(slide, "Mature chatbot projects separate UI, orchestration, LLM reasoning, API tools, and governance.")


def add_input_session(slide) -> None:
    add_step_slide(
        slide,
        "Step 1 - Capture User Input",
        "Create or resume a secure quote chat session.",
        ["Chat UI / Teams / Web", "SSO or Entra ID", "Session store", "Conversation ownership check", "Audit logging"],
        ["POST /chat/sessions", "GET /chat/sessions/{id}", "POST /chat/messages"],
        "Owner ask: IT must confirm approved channel, login model, and chat history retention.",
    )


def add_intent_extraction(slide) -> None:
    add_step_slide(
        slide,
        "Step 2 - Intent and Field Extraction",
        "Use the LLM to parse user language into structured quote fields.",
        ["LLM endpoint", "Intent taxonomy", "Entity schema", "Strict function schema", "Confidence fallback"],
        ["POST /nlp/parse"],
        "Owner ask: AI engineer needs 20-50 real sales questions from PLM / Sales for prompt tests.",
    )


def add_workflow_reasoning(slide) -> None:
    add_header(slide, "Step 3 - Missing-Field Reasoning")
    add_table(
        slide,
        ["Validation goal", "Required fields", "Endpoint"],
        [
            ["Region", "product_ids, region", "/workflow/required-fields"],
            ["System", "family, acquisition, stand, wallstand, table", "/workflow/{system_family}"],
            ["Detector / Grid", "grid_id, position, detector", "/workflow/next-question"],
            ["Generator / Tube", "generator, tube_spec, category", "/workflow/required-fields"],
        ],
        0.55,
        1.35,
        12.25,
        3.35,
        20,
    )
    add_process_row(slide, 0.82, 5.2, ["Extracted fields", "Required matrix", "Missing fields", "Clarifying question"], 2.75, 0.58)
    add_callout(slide, "Owner ask: PLM must confirm mandatory fields for FMT and OTC workflows.")


def add_grounding(slide) -> None:
    add_header(slide, "Step 4 - Product and Rule Grounding")
    add_process_row(slide, 0.75, 1.35, ["User wording", "Product search", "Source evidence", "Grounded result"], 2.75, 0.68)
    add_table(
        slide,
        ["Endpoint", "Purpose", "Owner question"],
        [
            ["GET /products/search", "Search ID, description, comments", "Who owns aliases?"],
            ["GET /products/{id}", "Product detail and source", "Who owns refresh?"],
            ["GET /products/{id}/constraints", "Product constraints", "Who approves wording?"],
            ["GET /rules/{rule_id}", "Rule evidence", "Who owns traceability?"],
        ],
        0.55,
        3.1,
        12.25,
        3.05,
        20,
    )
    add_callout(slide, "Urgent: confirm whether product grounding is file-based, SQL-based, or search-indexed.")


def add_quote_context(slide) -> None:
    add_header(slide, "Step 5 - Quote Context Contract")
    add_rect(slide, 0.55, 1.25, 6.1, 5.15, WHITE, BORDER)
    add_text(slide, "Required JSON fields", 0.88, 1.58, 5.2, 0.35, 26, ORANGE, bold=True)
    add_bullets(
        slide,
        [
            "session_id, snapshot_version, rule_set_version",
            "system_family, region, acquisition_type",
            "product_ids and configuration object",
            "stand, wallstand, table, grid, detector",
            "generator, tube_spec, spec_category",
        ],
        0.88,
        2.15,
        5.15,
        3.0,
        20,
    )
    add_rect(slide, 6.95, 1.25, 5.85, 5.15, LIGHT_ORANGE, ORANGE)
    add_text(slide, "Endpoint contract", 7.28, 1.58, 4.7, 0.35, 26, DEEP_ORANGE, bold=True)
    add_vertical_flow(slide, 7.3, 2.15, ["POST /quote-context", "PATCH /quote-context/{id}", "GET /quote-context/{id}"], 4.65, 0.62, 20)
    add_text(slide, "This contract prevents loose natural language from entering validation.", 7.3, 4.75, 4.75, 0.55, 20, BODY)
    add_callout(slide, "Owner ask: backend/API owner and rule engine owner must agree on this schema.")


def add_validation(slide) -> None:
    add_header(slide, "Step 6 - Deterministic Rule Validation")
    add_process_row(slide, 0.75, 1.35, ["Quote context", "POST /validation/check", "ValidationResult", "Bot explanation"], 2.75, 0.68)
    add_table(
        slide,
        ["Rule category", "Current behavior", "Required inputs"],
        [
            ["Region limit", "Blocks outside allowed regions", "product_ids + region"],
            ["System compatibility", "Blocks unsupported combos", "family + acquisition + equipment"],
            ["Detector / Grid", "Blocks unsupported pairing", "grid + position + detector"],
            ["Generator / Tube", "Returns specs or invalid", "generator + tube spec"],
        ],
        0.55,
        3.1,
        12.25,
        3.05,
        20,
    )
    add_callout(slide, "Critical rule: LLM explains results; rule engine decides results.")


def add_explanation(slide) -> None:
    add_step_slide(
        slide,
        "Step 7 - Explain and Recommend",
        "Convert deterministic validation results into clear sales-facing guidance.",
        ["Issue-code templates", "Source citation", "Next-action logic", "Sales wording review", "Escalation policy"],
        ["POST /explanations", "POST /recommendations"],
        "Owner ask: PLM and sales must approve wording for invalid / incomplete / warning responses.",
    )


def add_sme_review(slide) -> None:
    add_header(slide, "Step 8 - SME Review and Rule Governance")
    add_process_row(slide, 0.6, 1.3, ["Candidate rule", "SME decision", "Normalized payload", "Merged rules", "Engine handler"], 2.25, 0.68)
    add_table(
        slide,
        ["Decision", "Meaning", "Endpoint / artifact"],
        [
            ["approve", "Rule can become deterministic", "POST /rules/merge"],
            ["reject", "Not a valid business rule", "review case"],
            ["info_only", "Explanation only", "review case"],
            ["split", "One row becomes multiple rules", "reviewed rules"],
        ],
        0.55,
        3.1,
        12.25,
        3.05,
        20,
    )
    add_callout(slide, "Urgent business ask: assign owners for 387 rules needing review.")


def add_endpoint_matrix(slide) -> None:
    add_header(slide, "Urgent Endpoint Matrix")
    add_table(
        slide,
        ["Priority", "Endpoint", "Owner", "Why urgent"],
        [
            ["1", "POST /validation/check", "Rule engine", "Make MVP callable"],
            ["2", "GET /products/search", "Data/API", "Ground user wording"],
            ["3", "POST /quote-context", "Backend", "Shared schema"],
            ["4", "POST /workflow/required-fields", "PLM + dev", "Ask right questions"],
            ["5", "POST /nlp/parse", "AI engineer", "Extract intent"],
            ["6", "POST /review/cases", "SME + dev", "Govern ambiguity"],
        ],
        0.45,
        1.25,
        12.45,
        5.25,
        20,
    )
    add_callout(slide, "Every endpoint needs an owner, schema, host, auth model, and test examples.")


def add_collaboration(slide) -> None:
    add_header(slide, "Collaboration Model")
    add_table(
        slide,
        ["Role", "Needed from them", "Output"],
        [
            ["Mentor / sponsor", "Confirm framing", "Approved pilot scope"],
            ["PLM", "Confirm workflow + rules", "Step map + decisions"],
            ["IT / API", "Auth, hosting, standards", "Endpoint model"],
            ["Data engineer", "Snapshot refresh", "Reliable data pipeline"],
            ["AI engineer", "Prompt + tool calling", "Intent parser"],
            ["QA / UAT", "Real quote cases", "Evidence"],
        ],
        0.55,
        1.25,
        12.25,
        5.25,
        20,
    )
    add_callout(slide, "This cannot be solved by coding alone; ownership and SME review are the blockers.")


def add_decisions(slide) -> None:
    add_header(slide, "Decisions Needed to Start Pilot")
    add_rect(slide, 0.55, 1.25, 5.95, 5.2, WHITE, BORDER)
    add_text(slide, "First build", 0.88, 1.58, 4.5, 0.35, 26, ORANGE, bold=True)
    add_numbered_list(slide, ["Expose validation API", "Expose product search", "Define QuoteContext schema", "Create simple orchestrator", "Use 10-20 real quote cases"], 0.88, 2.1, 5.1, 3.4, 21)
    add_rect(slide, 6.85, 1.25, 5.95, 5.2, LIGHT_ORANGE, ORANGE)
    add_text(slide, "Company decisions", 7.18, 1.58, 4.5, 0.35, 26, DEEP_ORANGE, bold=True)
    add_bullets(slide, ["Approved chatbot channel", "Approved LLM platform", "API ownership", "Data source of truth", "SME rule-review owners", "UAT acceptance criteria"], 7.18, 2.1, 5.0, 3.4, 20)
    add_callout(slide, "Meeting close: confirm endpoint owners and first pilot scenario.")


def add_step_slide(slide, title: str, summary: str, support: list[str], endpoints: list[str], callout: str) -> None:
    add_header(slide, title)
    add_rect(slide, 0.55, 1.25, 5.95, 5.15, WHITE, BORDER)
    add_text(slide, "What happens", 0.88, 1.58, 4.5, 0.35, 26, ORANGE, bold=True)
    add_text(slide, summary, 0.88, 2.12, 5.0, 0.75, 22, TITLE, bold=True)
    add_text(slide, "Required support", 0.88, 3.32, 3.0, 0.3, 24, DEEP_ORANGE, bold=True)
    add_bullets(slide, support, 0.88, 3.78, 5.0, 2.0, 20)
    add_rect(slide, 6.85, 1.25, 5.95, 5.15, LIGHT_ORANGE, ORANGE)
    add_text(slide, "Endpoint examples", 7.18, 1.58, 4.5, 0.35, 26, DEEP_ORANGE, bold=True)
    add_vertical_flow(slide, 7.18, 2.18, endpoints, 4.9, 0.62, 20)
    add_text(slide, "Owner question", 7.18, 4.92, 2.2, 0.3, 24, TITLE, bold=True)
    add_text(slide, "Who can build, host, secure, and maintain this support?", 7.18, 5.34, 4.85, 0.45, 20, BODY)
    add_callout(slide, callout)


def add_header(slide, title: str) -> None:
    add_text(slide, title, 0.55, 0.34, 11.9, 0.55, 38, TITLE, bold=True)
    add_line(slide, 0.58, 1.02, 1.5, ORANGE)


def add_table(slide, headers: list[str], rows: list[list[str]], left: float, top: float, width: float, height: float, font_size: int) -> None:
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table
    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ORANGE
        format_cell(cell, WHITE, font_size, True)
    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_index % 2 else CARD
            format_cell(cell, BODY, font_size, False)


def add_process_row(slide, left: float, top: float, labels: list[str], box_width: float, box_height: float) -> None:
    for index, label in enumerate(labels):
        x = left + index * (box_width + 0.35)
        flow_box(slide, label, x, top, box_width, box_height, 20)
        if index < len(labels) - 1:
            add_text(slide, ">", x + box_width + 0.1, top + box_height / 2 - 0.1, 0.2, 0.2, 20, DEEP_ORANGE, bold=True)


def add_vertical_flow(slide, left: float, top: float, labels: list[str], box_w: float, box_h: float, font_size: int) -> None:
    for index, label in enumerate(labels):
        y = top + index * (box_h + 0.15)
        flow_box(slide, label, left, y, box_w, box_h, font_size)
        if index < len(labels) - 1:
            add_text(slide, "v", left + box_w / 2 - 0.05, y + box_h + 0.01, 0.2, 0.16, 20, DEEP_ORANGE, bold=True)


def flow_box(slide, text: str, left: float, top: float, width: float, height: float, font_size: int) -> None:
    add_rect(slide, left, top, width, height, WHITE, ORANGE)
    add_center_text(slide, text, left + 0.08, top + height / 2 - 0.13, width - 0.16, 0.22, font_size, TITLE, bold=True)


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float) -> None:
    if abs(x1 - x2) >= abs(y1 - y2):
        left = min(x1, x2)
        add_line(slide, left, y1, abs(x2 - x1), ORANGE)
        add_text(slide, ">" if x2 > x1 else "<", x2 - 0.06, y2 - 0.1, 0.18, 0.16, 20, ORANGE, bold=True)
    else:
        top = min(y1, y2)
        add_line(slide, x1, top, abs(y2 - y1), ORANGE, vertical=True)
        add_text(slide, "v" if y2 > y1 else "^", x2 - 0.07, y2 - 0.08, 0.16, 0.16, 20, ORANGE, bold=True)


def add_badge_row(slide, left: float, top: float, labels: list[str]) -> None:
    x = left
    for label in labels:
        width = max(1.6, len(label) * 0.1)
        add_rect(slide, x, top, width, 0.38, LIGHT_ORANGE, ORANGE)
        add_center_text(slide, label, x + 0.05, top + 0.04, width - 0.1, 0.22, 20, TITLE, bold=True)
        x += width + 0.22


def add_numbered_list(slide, items: list[str], left: float, top: float, width: float, height: float, size: int) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    for index, item in enumerate(items, start=1):
        paragraph = frame.paragraphs[0] if index == 1 else frame.add_paragraph()
        paragraph.text = f"{index}. {item}"
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = BODY
        paragraph.space_after = Pt(8)


def add_bullets(slide, bullets: list[str], left: float, top: float, width: float, height: float, size: int) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = BODY
        paragraph.space_after = Pt(6)


def add_text(slide, text: str, left: float, top: float, width: float, height: float, size: int, color: RGBColor, bold: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.line_spacing = 1.02


def add_center_text(slide, text: str, left: float, top: float, width: float, height: float, size: int, color: RGBColor, bold: bool = False) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def add_rect(slide, left: float, top: float, width: float, height: float, fill: RGBColor, line: RGBColor) -> None:
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)


def add_line(slide, left: float, top: float, length: float, color: RGBColor, vertical: bool = False) -> None:
    if vertical:
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(0.03), Inches(length))
    else:
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(length), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def add_callout(slide, text: str) -> None:
    add_rect(slide, 0.55, 6.55, 12.2, 0.45, LIGHT_ORANGE, ORANGE)
    add_text(slide, text, 0.82, 6.62, 11.5, 0.24, 20, TITLE, bold=True)


def format_cell(cell, color: RGBColor, size: int, bold: bool) -> None:
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color


def add_page_numbers(deck: Presentation) -> None:
    total = len(deck.slides)
    for index, slide in enumerate(deck.slides, start=1):
        add_text(slide, f"p.{index}/{total}", 11.45, 6.98, 1.25, 0.24, 20, BODY)


if __name__ == "__main__":
    main()